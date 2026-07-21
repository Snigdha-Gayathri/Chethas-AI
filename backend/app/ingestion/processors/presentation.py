from __future__ import annotations
from .base import BaseProcessor, ProcessedDocument
import logging

logger = logging.getLogger(__name__)

class PresentationProcessor(BaseProcessor):
    supported_types = ["application/vnd.openxmlformats-officedocument.presentationml.presentation"]
    
    async def process(self, file_path: str, metadata: dict = None) -> ProcessedDocument:
        metadata = metadata or {}
        text_content = ""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_text.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
            text_content = "\n\n".join(slides_text)
        except ImportError:
            logger.warning("python-pptx not installed.")
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            
        return ProcessedDocument(text_content=text_content, metadata=metadata)
